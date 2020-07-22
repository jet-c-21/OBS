import json

from pptx import Presentation


class PPT2JSON:
    def __init__(self, ppt_path: str):
        self.ppt_path = ppt_path
        self.ppt_name = ppt_path.split('/')[-1].replace('.pptx', '').strip()
        self.result = {'path': self.ppt_path,
                       'name': self.ppt_name,
                       'slides': list()}
        self.ppt_json = None

        self.__prs = None
        self.__load_ppt()
        self.__parse()

    def __load_ppt(self):
        self.__prs = Presentation(self.ppt_path)

    def __parse(self):
        for index, slide in enumerate(self.__prs.slides, start=1):
            temp = {'page': index}
            title = ''
            content = ''
            text = ''

            for el in slide.shapes:
                if hasattr(el, 'text'):
                    text += el.text
                    if ('標題' in el.name) or ('title' in el.name):
                        title += el.text
                    else:
                        content += el.text

            temp['title'] = title
            temp['content'] = content
            temp['text'] = text
            self.result['slides'].append(temp)

        self.ppt_json = json.dumps(self.result)

    def get_json(self) -> str:
        return self.ppt_json


# path = 'DATA/1.pptx'
#
# xx = PPT2JSON(path).get_json()
#
# print(json.loads(xx))
