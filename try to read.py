import os
from pptx import Presentation
from pprint import pprint as pp

path = 'DATA/人文素養-2.pptx'
prs = Presentation(path)

for index, slide in enumerate(prs.slides, start=1):
    print(f'第{index}頁')
    pp(dir(slide))
    for shape in slide.shapes:
        # print(dir(shape))
        # pp(dir(shape))
        if hasattr(shape, 'text'):
            print(shape.text, '**', shape.name)

    print('=' * 30)
